from pptx import Presentation
import os, hashlib

os.makedirs('static/images/energy', exist_ok=True)
saved = []

files = [
    ('ref/StarKernal_Smart_Energy_Solutions_EN.pptx', 'en'),
    ('ref/\u661f\u6838\u52a8\u529b_\u667a\u6167\u80fd\u6e90\u89e3\u51b3\u65b9\u6848v1.1.pptx', 'zh')
]

for fname, label in files:
    try:
        prs = Presentation(fname)
        for slide_num, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.shape_type == 13:
                    img = shape.image
                    ext = img.ext
                    h = hashlib.md5(img.blob).hexdigest()[:8]
                    out = 'static/images/energy/%s_slide%02d_%s.%s' % (label, slide_num, h, ext)
                    if not os.path.exists(out):
                        with open(out, 'wb') as f:
                            f.write(img.blob)
                        saved.append(out)
        print('Done: %s' % fname)
    except Exception as e:
        print('Error %s: %s' % (fname, e))

print('Extracted %d images:' % len(saved))
for s in saved:
    size = os.path.getsize(s)
    print('  %s (%dKB)' % (os.path.basename(s), size // 1024))
